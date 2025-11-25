"""
Import Service - Import content from various sources (PDF, PPTX, URL, Zoom)
"""

from typing import Optional, Dict, Any
import io

class ImportService:
    """Service for importing content from various sources"""
    
    @staticmethod
    async def import_from_pdf(file_content: bytes) -> Dict[str, Any]:
        """Import content from PDF file"""
        try:
            try:
                import PyPDF2
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                
                # Extract text from all pages
                pages = []
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    pages.append({
                        "page_number": page_num + 1,
                        "content": text,
                        "type": "text"
                    })
                
                return {
                    "title": "Imported from PDF",
                    "total_pages": len(pdf_reader.pages),
                    "pages": pages,
                    "format": "pdf"
                }
            except ImportError:
                raise Exception("PyPDF2 not installed. Run: pip install PyPDF2")
            
        except Exception as e:
            raise Exception(f"Failed to import PDF: {str(e)}")
    
    @staticmethod
    async def import_from_pptx(file_content: bytes) -> Dict[str, Any]:
        """Import content from PowerPoint file"""
        try:
            try:
                from pptx import Presentation as PPTXPresentation
                prs = PPTXPresentation(io.BytesIO(file_content))
                
                slides = []
                for slide_num, slide in enumerate(prs.slides):
                    slide_data = {
                        "slide_number": slide_num + 1,
                        "title": "",
                        "content": [],
                        "notes": ""
                    }
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            if shape.text and not slide_data["title"]:
                                slide_data["title"] = shape.text
                            else:
                                slide_data["content"].append({
                                    "type": "text",
                                    "content": shape.text
                                })
                    
                    # Extract notes
                    if slide.has_notes_slide:
                        notes_frame = slide.notes_slide.notes_text_frame
                        if notes_frame:
                            slide_data["notes"] = notes_frame.text
                    
                    slides.append(slide_data)
                
                return {
                    "title": "Imported from PowerPoint",
                    "total_slides": len(prs.slides),
                    "slides": slides,
                    "format": "pptx"
                }
            except ImportError:
                raise Exception("python-pptx not installed. Run: pip install python-pptx")
            
        except Exception as e:
            raise Exception(f"Failed to import PPTX: {str(e)}")
    
    @staticmethod
    async def import_from_url(url: str) -> Dict[str, Any]:
        """Import content from URL/webpage"""
        try:
            try:
                import requests
                from bs4 import BeautifulSoup
                from urllib.parse import urljoin
                
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                }
                response = requests.get(url, headers=headers, timeout=10)
                response.raise_for_status()
                
                soup = BeautifulSoup(response.content, 'html.parser')
                
                # Extract title
                title = soup.find('title')
                title_text = title.string if title else "Imported from Web"
                
                # Extract meta description
                meta_desc = soup.find('meta', attrs={'name': 'description'})
                description = meta_desc['content'] if meta_desc else ""
                
                # Extract main content
                main_content = soup.find('main') or soup.find('article') or soup.find('body')
                
                # Extract paragraphs and headings
                sections = []
                if main_content:
                    for element in main_content.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol']):
                        if element.name.startswith('h'):
                            sections.append({
                                "type": "heading",
                                "level": int(element.name[1]),
                                "content": element.get_text(strip=True)
                            })
                        elif element.name == 'p':
                            text = element.get_text(strip=True)
                            if text:
                                sections.append({
                                    "type": "paragraph",
                                    "content": text
                                })
                        elif element.name in ['ul', 'ol']:
                            items = [li.get_text(strip=True) for li in element.find_all('li')]
                            if items:
                                sections.append({
                                    "type": "list",
                                    "ordered": element.name == 'ol',
                                    "items": items
                                })
                
                # Extract images
                images = []
                for img in soup.find_all('img', src=True):
                    img_url = img['src']
                    if not img_url.startswith('http'):
                        img_url = urljoin(url, img_url)
                    
                    images.append({
                        "url": img_url,
                        "alt": img.get('alt', '')
                    })
                
                return {
                    "title": title_text,
                    "description": description,
                    "url": url,
                    "sections": sections,
                    "images": images[:10],  # Limit to 10 images
                    "format": "url"
                }
            except ImportError:
                raise Exception("Required libraries not installed. Run: pip install requests beautifulsoup4")
            
        except Exception as e:
            raise Exception(f"Failed to import from URL: {str(e)}")
    
    @staticmethod
    async def import_zoom_transcript(transcript_text: str) -> Dict[str, Any]:
        """Import and structure Zoom meeting transcript"""
        try:
            # Parse transcript (format: "Timestamp Speaker: Text")
            lines = transcript_text.split('\n')
            
            sections = []
            current_speaker = None
            current_content = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Try to parse speaker format
                if ':' in line:
                    parts = line.split(':', 1)
                    if len(parts) == 2:
                        # Check if this is a new speaker
                        potential_speaker = parts[0].strip()
                        if current_speaker != potential_speaker:
                            # Save previous speaker's content
                            if current_speaker and current_content:
                                sections.append({
                                    "type": "speaker",
                                    "speaker": current_speaker,
                                    "content": " ".join(current_content)
                                })
                            
                            current_speaker = potential_speaker
                            current_content = [parts[1].strip()]
                        else:
                            current_content.append(parts[1].strip())
                else:
                    current_content.append(line)
            
            # Add last speaker's content
            if current_speaker and current_content:
                sections.append({
                    "type": "speaker",
                    "speaker": current_speaker,
                    "content": " ".join(current_content)
                })
            
            return {
                "title": "Meeting Transcript",
                "sections": sections,
                "format": "zoom_transcript"
            }
            
        except Exception as e:
            raise Exception(f"Failed to import Zoom transcript: {str(e)}")
