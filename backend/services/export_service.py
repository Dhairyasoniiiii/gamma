"""
Export Service
Handles exporting presentations to various formats (PDF, PPTX, PNG, HTML, Markdown)
"""
import os
import json
from typing import Dict, List, Optional
from datetime import datetime
from io import BytesIO
import base64

# PDF export
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.lib.colors import HexColor
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# PowerPoint export
try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Image export
try:
    from PIL import Image as PILImage
    from PIL import ImageDraw, ImageFont
    IMAGE_AVAILABLE = True
except ImportError:
    IMAGE_AVAILABLE = False


class ExportService:
    """Service for exporting presentations to various formats"""
    
    def __init__(self):
        self.temp_dir = "/tmp/gamma_exports"
        os.makedirs(self.temp_dir, exist_ok=True)
    
    # ========== PDF Export ==========
    
    def export_to_pdf(
        self,
        presentation: dict,
        theme: Optional[dict] = None,
        output_path: Optional[str] = None
    ) -> str:
        """
        Export presentation to PDF
        
        Args:
            presentation: Presentation data with cards
            theme: Optional theme data
            output_path: Optional custom output path
            
        Returns:
            Path to generated PDF file
        """
        if not PDF_AVAILABLE:
            raise ImportError("ReportLab not installed. Run: pip install reportlab")
        
        # Generate filename
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(self.temp_dir, f"presentation_{timestamp}.pdf")
        
        # Create PDF document
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles based on theme
        if theme and theme.get("colors"):
            primary_color = theme["colors"].get("primary", "#000000")
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=HexColor(primary_color),
                alignment=TA_CENTER,
                spaceAfter=20
            )
        else:
            title_style = styles['Heading1']
        
        # Add title
        title = presentation.get("title", "Presentation")
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.5 * inch))
        
        # Add cards
        cards = presentation.get("content", {}).get("cards", [])
        for i, card in enumerate(cards):
            # Add card content based on type
            card_type = card.get("type", "text")
            
            if card_type == "title":
                # Title card
                if card.get("title"):
                    story.append(Paragraph(card["title"], styles['Heading1']))
                if card.get("subtitle"):
                    story.append(Paragraph(card["subtitle"], styles['Heading2']))
            
            elif card_type == "text":
                # Text card
                if card.get("title"):
                    story.append(Paragraph(card["title"], styles['Heading2']))
                if card.get("content"):
                    story.append(Paragraph(card["content"], styles['BodyText']))
            
            elif card_type == "list":
                # List card
                if card.get("title"):
                    story.append(Paragraph(card["title"], styles['Heading2']))
                items = card.get("items", [])
                for item in items:
                    story.append(Paragraph(f"• {item}", styles['BodyText']))
            
            elif card_type == "quote":
                # Quote card
                quote_style = ParagraphStyle(
                    'Quote',
                    parent=styles['Italic'],
                    fontSize=14,
                    leftIndent=20,
                    rightIndent=20
                )
                if card.get("content"):
                    story.append(Paragraph(f'"{card["content"]}"', quote_style))
                if card.get("author"):
                    story.append(Paragraph(f"— {card['author']}", styles['Normal']))
            
            # Add page break between cards (except last one)
            if i < len(cards) - 1:
                story.append(PageBreak())
            else:
                story.append(Spacer(1, 0.3 * inch))
        
        # Build PDF
        doc.build(story)
        return output_path
    
    # ========== PowerPoint Export ==========
    
    def export_to_pptx(
        self,
        presentation: dict,
        theme: Optional[dict] = None,
        output_path: Optional[str] = None
    ) -> str:
        """
        Export presentation to PowerPoint (PPTX)
        
        Args:
            presentation: Presentation data with cards
            theme: Optional theme data
            output_path: Optional custom output path
            
        Returns:
            Path to generated PPTX file
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")
        
        # Generate filename
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(self.temp_dir, f"presentation_{timestamp}.pptx")
        
        # Create PowerPoint presentation
        prs = PPTXPresentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add cards as slides
        cards = presentation.get("content", {}).get("cards", [])
        for card in cards:
            card_type = card.get("type", "text")
            
            # Choose slide layout
            if card_type == "title":
                slide_layout = prs.slide_layouts[0]  # Title slide
            else:
                slide_layout = prs.slide_layouts[1]  # Title and content
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Add content based on card type
            if card_type == "title":
                # Title slide
                if hasattr(slide.shapes, 'title'):
                    slide.shapes.title.text = card.get("title", "")
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = card.get("subtitle", "")
            
            elif card_type == "text":
                # Text slide
                if hasattr(slide.shapes, 'title'):
                    slide.shapes.title.text = card.get("title", "")
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.text = card.get("content", "")
            
            elif card_type == "list":
                # List slide
                if hasattr(slide.shapes, 'title'):
                    slide.shapes.title.text = card.get("title", "")
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()
                    
                    items = card.get("items", [])
                    for item in items:
                        p = text_frame.add_paragraph()
                        p.text = item
                        p.level = 0
            
            elif card_type == "quote":
                # Quote slide
                if hasattr(slide.shapes, 'title'):
                    slide.shapes.title.text = "Quote"
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    content = card.get("content", "")
                    author = card.get("author", "")
                    text_frame.text = f'"{content}"\n\n— {author}' if author else f'"{content}"'
        
        # Save presentation
        prs.save(output_path)
        return output_path
    
    # ========== HTML Export ==========
    
    def export_to_html(
        self,
        presentation: dict,
        theme: Optional[dict] = None,
        output_path: Optional[str] = None
    ) -> str:
        """
        Export presentation to HTML
        
        Args:
            presentation: Presentation data with cards
            theme: Optional theme data
            output_path: Optional custom output path
            
        Returns:
            Path to generated HTML file
        """
        # Generate filename
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(self.temp_dir, f"presentation_{timestamp}.html")
        
        # Get theme colors
        if theme and theme.get("colors"):
            primary = theme["colors"].get("primary", "#2563eb")
            background = theme["colors"].get("background", "#ffffff")
            text = theme["colors"].get("text", "#1f2937")
        else:
            primary = "#2563eb"
            background = "#ffffff"
            text = "#1f2937"
        
        # Build HTML
        html = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{presentation.get('title', 'Presentation')}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif;
            background-color: {background};
            color: {text};
            line-height: 1.6;
        }}
        
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            padding: 40px 20px;
        }}
        
        .slide {{
            background: white;
            border-radius: 12px;
            padding: 60px;
            margin-bottom: 40px;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
            min-height: 500px;
            display: flex;
            flex-direction: column;
            justify-content: center;
        }}
        
        .slide-title {{
            font-size: 48px;
            font-weight: bold;
            color: {primary};
            margin-bottom: 20px;
        }}
        
        .slide-subtitle {{
            font-size: 24px;
            color: #6b7280;
            margin-bottom: 20px;
        }}
        
        .slide-content {{
            font-size: 18px;
            margin-bottom: 20px;
        }}
        
        .slide-list {{
            list-style: none;
            padding-left: 0;
        }}
        
        .slide-list li {{
            font-size: 18px;
            margin-bottom: 12px;
            padding-left: 30px;
            position: relative;
        }}
        
        .slide-list li:before {{
            content: "•";
            color: {primary};
            font-size: 24px;
            position: absolute;
            left: 0;
        }}
        
        .slide-quote {{
            font-size: 24px;
            font-style: italic;
            text-align: center;
            padding: 40px;
            border-left: 4px solid {primary};
            margin: 20px 0;
        }}
        
        .slide-author {{
            text-align: right;
            color: #6b7280;
            margin-top: 20px;
        }}
        
        @media print {{
            .slide {{
                page-break-after: always;
            }}
        }}
    </style>
</head>
<body>
    <div class="container">
"""
        
        # Add cards
        cards = presentation.get("content", {}).get("cards", [])
        for card in cards:
            card_type = card.get("type", "text")
            html += '        <div class="slide">\n'
            
            if card_type == "title":
                if card.get("title"):
                    html += f'            <h1 class="slide-title">{card["title"]}</h1>\n'
                if card.get("subtitle"):
                    html += f'            <p class="slide-subtitle">{card["subtitle"]}</p>\n'
            
            elif card_type == "text":
                if card.get("title"):
                    html += f'            <h2 class="slide-title">{card["title"]}</h2>\n'
                if card.get("content"):
                    html += f'            <p class="slide-content">{card["content"]}</p>\n'
            
            elif card_type == "list":
                if card.get("title"):
                    html += f'            <h2 class="slide-title">{card["title"]}</h2>\n'
                items = card.get("items", [])
                if items:
                    html += '            <ul class="slide-list">\n'
                    for item in items:
                        html += f'                <li>{item}</li>\n'
                    html += '            </ul>\n'
            
            elif card_type == "quote":
                html += '            <blockquote class="slide-quote">\n'
                if card.get("content"):
                    html += f'                "{card["content"]}"\n'
                if card.get("author"):
                    html += f'                <div class="slide-author">— {card["author"]}</div>\n'
                html += '            </blockquote>\n'
            
            html += '        </div>\n'
        
        html += """
    </div>
</body>
</html>
"""
        
        # Write to file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
        
        return output_path
    
    # ========== Markdown Export ==========
    
    def export_to_markdown(
        self,
        presentation: dict,
        output_path: Optional[str] = None
    ) -> str:
        """
        Export presentation to Markdown
        
        Args:
            presentation: Presentation data with cards
            output_path: Optional custom output path
            
        Returns:
            Path to generated Markdown file
        """
        # Generate filename
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(self.temp_dir, f"presentation_{timestamp}.md")
        
        # Build Markdown
        md = f"# {presentation.get('title', 'Presentation')}\n\n"
        md += "---\n\n"
        
        # Add cards
        cards = presentation.get("content", {}).get("cards", [])
        for i, card in enumerate(cards):
            card_type = card.get("type", "text")
            
            if card_type == "title":
                if card.get("title"):
                    md += f"# {card['title']}\n\n"
                if card.get("subtitle"):
                    md += f"## {card['subtitle']}\n\n"
            
            elif card_type == "text":
                if card.get("title"):
                    md += f"## {card['title']}\n\n"
                if card.get("content"):
                    md += f"{card['content']}\n\n"
            
            elif card_type == "list":
                if card.get("title"):
                    md += f"## {card['title']}\n\n"
                items = card.get("items", [])
                for item in items:
                    md += f"- {item}\n"
                md += "\n"
            
            elif card_type == "quote":
                if card.get("content"):
                    md += f"> {card['content']}\n"
                if card.get("author"):
                    md += f">\n> — {card['author']}\n"
                md += "\n"
            
            # Add separator between cards
            if i < len(cards) - 1:
                md += "---\n\n"
        
        # Write to file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(md)
        
        return output_path


# Singleton instance
export_service = ExportService()
